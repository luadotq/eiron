# EIRON - Extensible Intelligent Retrieval Open-source Navigator
# GNU GPL V3 License
# Windows version
# Semantic search included
import os
import re
import time
import pickle
import zlib
import heapq
import threading
import csv
from dataclasses import dataclass
from typing import Dict, List, Tuple, Optional
from concurrent.futures import ThreadPoolExecutor
from queue import Queue, Empty
import psutil
import warnings
from collections import defaultdict
import numpy as np
from bs4 import BeautifulSoup
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import pandas as pd

# Configuration
MAX_RAM_BYTES = 12 * 1024 * 1024 * 1024  # 12 GB
CHUNK_SIZE = 250_000
MAX_DOCS_IN_MEMORY = 3000
MIN_TOKEN_LEN = 3
INDEX_FILE = "search_index.pkl"
TMP_INDEX_FILE = "search_index_tmp.pkl"
CONTEXT_WINDOW = 0  # Context window size
SUPPORTED_EXTENSIONS = {
    '.txt', '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.csv', 
    '.html', '.htm', '.pptx', '.ppt', '.odt', '.rtf'
}

@dataclass
class DocumentInfo:
    path: str
    excerpt: str
    size: int
    last_modified: float
    chunk_id: int = 0
    full_text: str = ""
    vector: Optional[np.ndarray] = None

class TextProcessor:
    
    @staticmethod
    #Text Normalize
    def preprocess_text(text: str) -> str:
        text = re.sub(r'\s+', ' ', text.lower())
        text = re.sub(r'[^\w\s\-]', '', text)
        return text
    #Split
    @staticmethod
    def split_sentences(text: str) -> List[str]:
        return re.split(r'(?<=[.!?])\s+', text)
    #Context extraction
    @staticmethod
    def extract_context(text: str, keywords: List[str]) -> str:
        sentences = TextProcessor.split_sentences(text)
        relevant = []
    
        for i, sentence in enumerate(sentences):
            if any(keyword in sentence for keyword in keywords):
                relevant.append(sentence)
            
                if CONTEXT_WINDOW >= 1:
                    if i > 0:
                        relevant.append(sentences[i-1])
                    if i < len(sentences)-1:
                        relevant.append(sentences[i+1])
    
        return ' '.join(relevant) if relevant else text[:500]

    @staticmethod
    def parse_file(filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == '.pdf':
                return TextProcessor._parse_pdf(filepath)
            elif ext == '.csv':
                return TextProcessor._parse_csv(filepath)
            elif ext in ('.html', '.htm'):
                return TextProcessor._parse_html(filepath)
            elif ext in ('.xlsx', '.xls'):
                return TextProcessor._parse_excel(filepath)
            elif ext == '.docx':
                return TextProcessor._parse_docx(filepath)
            elif ext == '.doc':
                return TextProcessor._parse_doc(filepath)
            elif ext == '.pptx':
                return TextProcessor._parse_pptx(filepath)
            elif ext == '.ppt':
                return TextProcessor._parse_ppt(filepath)
            elif ext == '.odt':
                return TextProcessor._parse_odt(filepath)
            elif ext == '.txt':
                return TextProcessor._parse_txt(filepath)
            elif ext == '.rtf':
                return TextProcessor._parse_rtf(filepath)
            else:
                print(f"[-] Unsupported file format: {ext}")
                return ""
        except Exception as e:
            print(f"[-] Error parsing {filepath}: {e}")
            return ""

    @staticmethod
    def _parse_pdf(filepath: str) -> str:
        from pdfminer.high_level import extract_text
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            return extract_text(filepath)

    @staticmethod
    def _parse_csv(filepath: str) -> str:
        content = []
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                content.append(' '.join(row))
        return '\n'.join(content)

    @staticmethod
    def _parse_html(filepath: str) -> str:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
            for script in soup(["script", "style"]):
                script.decompose()
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            return '\n'.join(chunk for chunk in chunks if chunk)

    @staticmethod
    def _parse_excel(filepath: str) -> str:
        content = []
        try:
            xls = pd.ExcelFile(filepath)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                content.append(f"=== Sheet: {sheet_name} ===")
                content.append(df.to_string())
        except Exception as e:
            print(f"[-] Error parsing Excel file {filepath}: {e}")
        return '\n'.join(content)

    @staticmethod
    def _parse_docx(filepath: str) -> str:
        from docx import Document
        doc = Document(filepath)
        return '\n'.join(p.text for p in doc.paragraphs if p.text)

    @staticmethod
    def _parse_doc(filepath: str) -> str:
        try:
            import subprocess
            result = subprocess.run(['antiword', filepath], capture_output=True, text=True)
            return result.stdout if result.returncode == 0 else ""
        except Exception as e:
            print(f"[-] Error parsing DOC file {filepath}: {e}")
            return ""

    @staticmethod
    def _parse_pptx(filepath: str) -> str:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

    @staticmethod
    def _parse_ppt(filepath: str) -> str:
        try:
            import comtypes.client
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            presentation = powerpoint.Presentations.Open(filepath)
            text = []
            for slide in presentation.Slides:
                for shape in slide.Shapes:
                    if hasattr(shape, "TextFrame"):
                        text.append(shape.TextFrame.TextRange.Text)
            presentation.Close()
            powerpoint.Quit()
            return '\n'.join(text)
        except Exception as e:
            print(f"[-] Error parsing PPT file {filepath}: {e}")
            return ""

    @staticmethod
    def _parse_odt(filepath: str) -> str:
        from odf.opendocument import load
        from odf.text import P
        doc = load(filepath)
        text = []
        for para in doc.getElementsByType(P):
            text.append(para.__str__())
        return '\n'.join(text)

    @staticmethod
    def _parse_rtf(filepath: str) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return rtf_to_text(f.read())
        except Exception as e:
            print(f"[-] Error parsing RTF file {filepath}: {e}")
            return ""

    @staticmethod
    def _parse_txt(filepath: str) -> str:
        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

class HybridSearchEngine:
    def __init__(self):
        self.inverted_index: Dict[str, Dict[Tuple[int, int], int]] = {}
        self.documents: Dict[int, DocumentInfo] = {}
        self.doc_path_to_id: Dict[str, int] = {}
        self.next_doc_id = 1
        self.tfidf = TfidfVectorizer(stop_words='english', max_features=5000)
        self.is_tfidf_trained = False
        self.index_queue = Queue(maxsize=5000)
        self.shutdown_flag = False
        self.lock = threading.Lock()
        self.last_flush = time.time()
        
        self._load_index()
        self._start_workers()

    def _start_workers(self):
        num_workers = min(8, (os.cpu_count() or 4))
        self.workers = []
        for _ in range(num_workers):
            t = threading.Thread(target=self._process_queue, daemon=True)
            t.start()
            self.workers.append(t)

        self.memory_manager = threading.Thread(target=self._manage_memory, daemon=True)
        self.memory_manager.start()
#Memory management
    def _manage_memory(self):
        while not self.shutdown_flag:
            try:
                mem = psutil.virtual_memory()
                used_percent = mem.used / MAX_RAM_BYTES

                if used_percent > 0.8 or len(self.documents) >= MAX_DOCS_IN_MEMORY:
                    self._flush_to_disk()

                time.sleep(1)
            except Exception as e:
                print(f"Memory error: {e}")
                break

    def _flush_to_disk(self):
        with self.lock:
            if not self.documents:
                return

            print(f"Saving {len(self.documents)} documents...")
            start = time.time()

            try:
                data = {
                    'inverted_index': self.inverted_index,
                    'documents': self.documents,
                    'doc_path_to_id': self.doc_path_to_id,
                    'next_doc_id': self.next_doc_id,
                    'tfidf_vocab': self.tfidf.vocabulary_ if self.is_tfidf_trained else None,
                    'tfidf_idf': self.tfidf.idf_ if self.is_tfidf_trained else None
                }

                with open(TMP_INDEX_FILE, 'wb') as f:
                    f.write(zlib.compress(pickle.dumps(data, protocol=pickle.HIGHEST_PROTOCOL)))

                os.replace(TMP_INDEX_FILE, INDEX_FILE)
                self._clear_memory()

                print(f"Saved in {time.time() - start:.2f}s")
            except Exception as e:
                print(f"Save failed: {e}")
                self._clear_memory()

    def _clear_memory(self):
        self.inverted_index.clear()
        self.documents.clear()
        self.doc_path_to_id.clear()

    def _load_index(self):
        if os.path.exists(INDEX_FILE):
            try:
                print("Loading index...")
                with open(INDEX_FILE, 'rb') as f:
                    data = pickle.loads(zlib.decompress(f.read()))
                
                self.inverted_index = data.get('inverted_index', {})
                self.documents = data.get('documents', {})
                self.doc_path_to_id = data.get('doc_path_to_id', {})
                self.next_doc_id = data.get('next_doc_id', 1)
                
                if data.get('tfidf_vocab') and data.get('tfidf_idf'):
                    self.tfidf.vocabulary_ = data['tfidf_vocab']
                    self.tfidf.idf_ = data['tfidf_idf']
                    self.is_tfidf_trained = True
                    self._update_document_vectors()
                
                print(f"Loaded {len(self.documents)} documents")
            except Exception as e:
                print(f"Load error: {e}")
                self.inverted_index = {}
                self.documents = {}
                self.doc_path_to_id = {}
                self.next_doc_id = 1
#Semantic model: Vectrors
    def _update_document_vectors(self):
        if not self.is_tfidf_trained or not self.documents:
            return
        
        try:
            texts = [doc.full_text for doc in self.documents.values() if doc.full_text]
            if texts:
                vectors = self.tfidf.transform(texts)
                for (doc_id, doc), vector in zip(self.documents.items(), vectors):
                    doc.vector = vector
        except Exception as e:
            print(f"[-] Error updating document vectors: {e}")

    def _process_queue(self):
        while not self.shutdown_flag:
            try:
                doc_info, tokens = self.index_queue.get(timeout=1)

                with self.lock:
                    doc_id = self.doc_path_to_id.get(doc_info.path, self.next_doc_id)
                    if doc_id == self.next_doc_id:
                        self.doc_path_to_id[doc_info.path] = doc_id
                        self.next_doc_id += 1

                    self.documents[doc_id] = doc_info

                    for token, count in tokens.items():
                        if token not in self.inverted_index:
                            self.inverted_index[token] = {}
                        key = (doc_id, doc_info.chunk_id)
                        self.inverted_index[token][key] = count

                    if len(self.documents) % 100 == 0:
                        self._train_tfidf()

            except Empty:
                continue
            except Exception as e:
                print(f"[-] Queue error: {e}")
#Semantic Model: TF-IDF Model training
    def _train_tfidf(self):
        texts = [doc.full_text for doc in self.documents.values() if doc.full_text]
        if texts:
            self.tfidf.fit(texts)
            self.is_tfidf_trained = True
            self._update_document_vectors()

    def _tokenize(self, text: str) -> Dict[str, int]:
        text = TextProcessor.preprocess_text(text)
        tokens = re.findall(r'\b[\w\-]{%d,}\b' % MIN_TOKEN_LEN, text)
        
        counts = defaultdict(int)
        for t in tokens:
            if 3 <= len(t) <= 25 and not t.isnumeric():
                counts[t] += 1
        return dict(counts)

    def index_file(self, filepath: str):
        try:
            filepath = os.path.abspath(filepath)
            ext = os.path.splitext(filepath)[1].lower()
            
            if ext not in SUPPORTED_EXTENSIONS:
                print(f"[-] Skipping unsupported file: {filepath}")
                return

            stat = os.stat(filepath)
            if filepath in self.doc_path_to_id:
                doc_id = self.doc_path_to_id[filepath]
                existing_doc = self.documents.get(doc_id)
                if existing_doc and existing_doc.last_modified >= stat.st_mtime:
                    print(f"File {os.path.basename(filepath)} already indexed (up to date)")
                    return

            full_text = TextProcessor.parse_file(filepath)
            if not full_text.strip():
                print(f"[-] Empty content in file: {filepath}")
                return

            excerpt = full_text[:300] + '...' if len(full_text) > 300 else full_text
            tokens = self._tokenize(full_text)

            doc_info = DocumentInfo(
                path=filepath,
                excerpt=excerpt,
                size=stat.st_size,
                last_modified=stat.st_mtime,
                full_text=full_text
            )

            self.index_queue.put((doc_info, tokens))
            print(f"Indexed {os.path.basename(filepath)}")

        except Exception as e:
            print(f"[-] Indexing error {filepath}: {e}")

    def index_directory(self, directory: str):
        if not os.path.isdir(directory):
            print(f"Invalid directory: {directory}")
            return

        print(f"Indexing {directory}...")
        start = time.time()
        file_count = 0

        with ThreadPoolExecutor(max_workers=8) as executor:
            futures = []
            for root, _, files in os.walk(directory):
                for file in files:
                    filepath = os.path.join(root, file)
                    if os.path.splitext(filepath)[1].lower() in SUPPORTED_EXTENSIONS:
                        futures.append(executor.submit(self.index_file, filepath))
                        file_count += 1

                        if file_count % 10 == 0:
                            print(f"Processed {file_count} files")

            for future in futures:
                try:
                    future.result()
                except Exception as e:
                    print(f"Error: {e}")

        self._flush_to_disk()
        print(f"Indexed {file_count} files in {time.time() - start:.1f}s")
#Search and semantic search
    def search(self, query: str, limit: int = 10) -> List[Dict]:
        if not query.strip():
            return []

        start_time = time.time()
        query = TextProcessor.preprocess_text(query)
        query_tokens = list(self._tokenize(query).keys())

        if not self.documents:
            self._load_index()

        # Keyword search
        keyword_scores = defaultdict(float)
        with self.lock:
            for token in query_tokens:
                if token in self.inverted_index:
                    for (doc_id, _), count in self.inverted_index[token].items():
                        keyword_scores[doc_id] += count * 1.5

        # Semantic model search
        semantic_scores = defaultdict(float)
        if self.is_tfidf_trained and query_tokens:
            query_vec = self.tfidf.transform([query])
            for doc_id, doc in self.documents.items():
                if doc.vector is not None:
                    similarity = cosine_similarity(query_vec, doc.vector.reshape(1, -1))[0][0]
                    semantic_scores[doc_id] += similarity * 2.0

        # Creating results
        combined_scores = defaultdict(float)
        all_doc_ids = set(keyword_scores.keys()).union(set(semantic_scores.keys()))
        for doc_id in all_doc_ids:
            combined_scores[doc_id] = keyword_scores.get(doc_id, 0) + semantic_scores.get(doc_id, 0)

        # Finishing results
        results = []
        for doc_id, score in heapq.nlargest(limit, combined_scores.items(), key=lambda x: x[1]):
            doc = self.documents.get(doc_id)
            if doc:
                context = TextProcessor.extract_context(doc.full_text, query_tokens)
                results.append({
                    'path': doc.path,
                    'excerpt': context,
                    'score': score,
                    'modified': time.strftime('%Y-%m-%d %H:%M', time.localtime(doc.last_modified)),
                    'size': f"{doc.size / 1024:.1f} KB"
                })

        print(f"Search completed in {time.time() - start_time:.3f}s")
        return sorted(results, key=lambda x: x['score'], reverse=True)[:limit]

    def shutdown(self):
        print("Shutting down...")
        self.shutdown_flag = True
        self._flush_to_disk()
        for t in self.workers:
            t.join(timeout=1)
        self.memory_manager.join(timeout=1)

if __name__ == "__main__":
    engine = HybridSearchEngine()
    try:
        while True:
            print("\n1. Index directory\n2. Search\n3. Exit")
            choice = input("> ").strip()

            if choice == "1":
                path = input("Directory path: ").strip()
                if os.path.isdir(path):
                    engine.index_directory(path)
                else:
                    print("[-] Invalid directory path")
            elif choice == "2":
                query = input("Search query: ").strip()
                results = engine.search(query)
                for i, res in enumerate(results, 1):
                    print(f"\n{i}. {res['path']} (score: {res['score']:.2f}, {res['size']}, modified: {res['modified']})")
                    print(f"   {res['excerpt']}")
            elif choice == "3":
                break
    except KeyboardInterrupt:
        print("\nInterrupted by user")
    finally:
        engine.shutdown()
